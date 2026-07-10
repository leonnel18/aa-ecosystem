"""
Generate AktivAsia Portal documentation:
  - docs/AktivAsia_Portal_User_Guide.docx
  - docs/AktivAsia_Portal_Presentation.pptx

Run: python tools/generate_docs.py
"""

from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

from doc_content import COVER, SECTIONS
from pptx_content import SLIDES

# Brand colors
BURGUNDY = RGBColor(0x82, 0x15, 0x45)
ORANGE   = RGBColor(0xE8, 0x6A, 0x1A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF4, 0xEE)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY   = RGBColor(0x6B, 0x6B, 0x6B)

BURGUNDY_PPT = PptRGB(0x82, 0x15, 0x45)
ORANGE_PPT   = PptRGB(0xE8, 0x6A, 0x1A)
WHITE_PPT    = PptRGB(0xFF, 0xFF, 0xFF)
LIGHT_BG_PPT = PptRGB(0xF5, 0xF4, 0xEE)
DARK_PPT     = PptRGB(0x1A, 0x1A, 0x1A)
MID_GREY_PPT = PptRGB(0x6B, 0x6B, 0x6B)

OUT_DIR      = Path(__file__).parent.parent / "docs"
SCREENS_DIR  = OUT_DIR / "screenshots"
OUT_DIR.mkdir(exist_ok=True)

def screen(name: str) -> str | None:
    """Return absolute path to a screenshot if it exists, else None."""
    p = SCREENS_DIR / f"{name}.png"
    return str(p) if p.exists() else None

# ─── WORD DOCUMENT ────────────────────────────────────────────────────────────

def set_cell_bg(cell, hex_color: str):
    """Set table cell background color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def add_paragraph_border_bottom(paragraph, color="821545"):
    """Add a bottom border to a paragraph (used under headings)."""
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "4")
    bottom.set(qn("w:space"), "4")
    bottom.set(qn("w:color"), color)
    pBdr.append(bottom)
    pPr.append(pBdr)


def style_run(run, bold=False, italic=False, size=None, color=None, font="Calibri"):
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    if size:
        run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def add_cover_page(doc):
    # Full burgundy background via a shaded table spanning page
    section = doc.sections[0]
    section.page_width  = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin   = Inches(1)
    section.right_margin  = Inches(1)
    section.top_margin    = Inches(1)
    section.bottom_margin = Inches(1)

    # Top accent bar
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    style_run(run, color=BURGUNDY, size=8)

    doc.add_paragraph()

    # Logo placeholder text
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("🌏  AktivAsia")
    style_run(run, bold=True, size=28, color=BURGUNDY, font="Calibri")

    doc.add_paragraph()

    # Title
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(COVER["title"])
    style_run(run, bold=True, size=32, color=BURGUNDY, font="Calibri")

    # Subtitle
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(COVER["subtitle"])
    style_run(run, size=16, color=ORANGE, font="Calibri")

    doc.add_paragraph()

    # Tagline box (shaded)
    table = doc.add_table(rows=1, cols=1)
    table.style = "Table Grid"
    cell = table.cell(0, 0)
    set_cell_bg(cell, "F5F4EE")
    cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = cell.paragraphs[0].add_run(COVER["tagline"])
    style_run(run, size=13, color=DARK_TEXT, font="Calibri")
    cell.paragraphs[0].paragraph_format.space_before = Pt(8)
    cell.paragraphs[0].paragraph_format.space_after  = Pt(8)

    doc.add_paragraph()
    doc.add_paragraph()
    doc.add_paragraph()

    # Author block
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(COVER["author"])
    style_run(run, bold=True, size=13, color=DARK_TEXT, font="Calibri")

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(COVER["role"])
    style_run(run, size=11, color=MID_GREY, font="Calibri")

    doc.add_paragraph()

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(f"{COVER['date']}  |  {COVER['version']}")
    style_run(run, size=10, color=MID_GREY, font="Calibri")

    # Bottom accent bar
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    style_run(run, color=BURGUNDY, size=8)

    doc.add_page_break()


def add_section(doc, section_data):
    heading = section_data["heading"]

    # Section heading
    p = doc.add_paragraph()
    run = p.add_run(heading)
    style_run(run, bold=True, size=16, color=BURGUNDY, font="Calibri")
    add_paragraph_border_bottom(p)
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after  = Pt(6)

    # Section intro
    if section_data.get("intro"):
        p = doc.add_paragraph()
        run = p.add_run(section_data["intro"])
        style_run(run, size=11, color=DARK_TEXT, font="Calibri")
        p.paragraph_format.space_after = Pt(10)

    # Find applicable screenshot rules for this section
    section_screens = {}
    for key, entries in DOCX_SCREENSHOTS.items():
        if key in heading:
            for sub_frag, img_name, caption in entries:
                section_screens.setdefault(sub_frag, []).append((img_name, caption))

    for sub in section_data.get("subsections", []):
        sub_title = sub["title"]

        # Subsection title
        p = doc.add_paragraph()
        run = p.add_run(sub_title)
        style_run(run, bold=True, size=12, color=ORANGE, font="Calibri")
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after  = Pt(4)

        # Subsection body — render line by line
        for line in sub["body"].split("\n"):
            p = doc.add_paragraph()
            run = p.add_run(line if line else " ")
            style_run(run, size=11, color=DARK_TEXT, font="Calibri")
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after  = Pt(2)
            if line.startswith("   ") or line.startswith("    "):
                p.paragraph_format.left_indent = Inches(0.3)

        # Inject screenshots after this subsection if any match
        for frag, entries in section_screens.items():
            if frag in sub_title:
                for img_name, caption in entries:
                    add_screenshot_docx(doc, screen(img_name), caption)

    doc.add_paragraph()


def add_screenshot_docx(doc, img_path: str, caption: str = "", width=Inches(5.5)):
    """Insert a screenshot image centered with an optional caption."""
    if not img_path:
        return
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run()
    run.add_picture(img_path, width=width)
    if caption:
        p2 = doc.add_paragraph()
        p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r2 = p2.add_run(f"↑  {caption}")
        style_run(r2, size=9, italic=True, color=MID_GREY, font="Calibri")
    doc.add_paragraph()


# Screenshot map: section heading fragment → list of (subsection_title_fragment, screenshot_name, caption)
DOCX_SCREENSHOTS = {
    "3. Navigating": [
        ("Step 2", "01_landing_page", "The Landing Page — showing all 5 country portals and grand totals"),
        ("Step 3", "02_portal_ph_overview", "The 4 tabs at the top of a country portal"),
    ],
    "4.1": [
        ("Overview Tab", "02_portal_ph_overview", "Overview Tab — Total Applicants, Graduates, Ongoing, and the Funnel"),
    ],
    "4.2": [
        ("Training Plan Tab", "03_portal_ph_training_plan", "Training Plan Tab — Trainings by type, upcoming, and completed"),
    ],
    "4.3": [
        ("Impact Tab", "04_portal_ph_impact", "Impact Tab — Likert skill scores and post-survey feedback"),
    ],
    "4.4": [
        ("Demographics Tab", "05_portal_ph_demographics", "Demographics Tab — Gender, age group, and region breakdown"),
    ],
    "5. How To: Manage": [
        ("5.1", "07_admin_panel", "The Admin Panel — 3 tabs: Selection, Attendance, Graduates"),
        ("5.2", "08_admin_selection_tab", "Selection Tab — review and select applicants"),
        ("5.3", "09_admin_attendance_tab", "Attendance Tab — mark who attended the training"),
        ("5.4", "10_admin_graduates_tab", "Graduates Tab — track post-survey completion and 6-Month Eval progress"),
    ],
    "6. How To: Create": [
        ("Step-by-Step", "11_create_training", "The Create Training form — fill in all fields then click Create Training"),
    ],
    "7. How To: Share": [
        ("The Application Form", "12_apply_ph", "The Philippines Application Form — what applicants see when they register"),
        ("The Post-Survey Form", "13_post_survey_ph", "The Post-Survey Form — sent to participants after training ends"),
    ],
    "2. What Is": [
        ("The 5 Portals", "06_backbone_overview", "The Backbone Portal — aggregate view across all countries"),
    ],
}


def build_docx():
    doc = Document()

    # Default style tweaks
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    add_cover_page(doc)

    # Table of contents header
    p = doc.add_paragraph()
    run = p.add_run("📋  Table of Contents")
    style_run(run, bold=True, size=14, color=BURGUNDY, font="Calibri")
    add_paragraph_border_bottom(p)
    p.paragraph_format.space_after = Pt(6)

    toc_items = [
        ("1.", "Welcome", "👋"),
        ("2.", "What Is the AktivAsia Portal?", "🌏"),
        ("3.", "Navigating the Portal", "🗺️"),
        ("4.", "How To: Read Your Reports", "📈"),
        ("5.", "How To: Manage a Training (Admin Panel)", "⚙️"),
        ("6.", "How To: Create a New Training", "➕"),
        ("7.", "How To: Share Application & Survey Forms", "📋"),
        ("8.", "Automated Reminders", "🔔"),
        ("9.", "Glossary", "📖"),
        ("10.", "Need Help?", "📞"),
    ]
    for num, title, emoji in toc_items:
        p = doc.add_paragraph()
        run = p.add_run(f"  {num}  {emoji}  {title}")
        style_run(run, size=11, color=DARK_TEXT, font="Calibri")
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after  = Pt(2)

    doc.add_page_break()

    for section_data in SECTIONS:
        add_section(doc, section_data)

    # Footer note
    p = doc.add_paragraph()
    run = p.add_run(
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━\n"
        f"AktivAsia Portal User Guide  |  {COVER['version']}  |  {COVER['date']}  |  "
        f"Prepared by {COVER['author']}, {COVER['role']}"
    )
    style_run(run, size=8, color=MID_GREY, font="Calibri")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    out = OUT_DIR / "AktivAsia_Portal_User_Guide.docx"
    doc.save(str(out))
    print(f"✅  Saved: {out}")


# ─── POWERPOINT ───────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Emu(12700)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 bold=False, italic=False, size=18, color=DARK_PPT,
                 align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.name   = font
    run.font.size   = PptPt(size)
    run.font.color.rgb = color
    return txBox


def add_footer(slide, text):
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), fill_color=BURGUNDY_PPT)
    add_text_box(
        slide, text,
        Inches(0.3), SLIDE_H - Inches(0.38),
        SLIDE_W - Inches(0.6), Inches(0.36),
        size=9, color=WHITE_PPT, align=PP_ALIGN.CENTER, font="Calibri"
    )


def add_slide_title(slide, title, y=Inches(0.35)):
    add_text_box(
        slide, title,
        Inches(0.5), y, SLIDE_W - Inches(1), Inches(0.7),
        bold=True, size=28, color=BURGUNDY_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )
    # Orange underline bar
    add_rect(slide, Inches(0.5), y + Inches(0.65), Inches(4), Inches(0.05), fill_color=ORANGE_PPT)


def build_cover(prs, data):
    slide = blank_slide(prs)

    # Full burgundy left panel
    add_rect(slide, 0, 0, Inches(5.5), SLIDE_H, fill_color=BURGUNDY_PPT)

    # Logo text on left panel
    add_text_box(
        slide, "🌏  AktivAsia",
        Inches(0.4), Inches(1.2), Inches(4.8), Inches(0.8),
        bold=True, size=32, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )

    # Main title on left
    add_text_box(
        slide, data["title"],
        Inches(0.4), Inches(2.2), Inches(4.8), Inches(1.4),
        bold=True, size=30, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )

    # Orange accent bar
    add_rect(slide, Inches(0.4), Inches(3.7), Inches(2), Inches(0.08), fill_color=ORANGE_PPT)

    # Subtitle
    add_text_box(
        slide, data["subtitle"],
        Inches(0.4), Inches(3.9), Inches(4.8), Inches(0.9),
        size=14, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )

    # Right side - light background
    add_rect(slide, Inches(5.5), 0, SLIDE_W - Inches(5.5), SLIDE_H, fill_color=LIGHT_BG_PPT)

    # Author info on right
    add_text_box(
        slide, data["author"],
        Inches(5.9), Inches(2.5), Inches(7), Inches(0.6),
        bold=True, size=20, color=BURGUNDY_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )
    add_text_box(
        slide, data["role"],
        Inches(5.9), Inches(3.1), Inches(7), Inches(0.5),
        size=13, color=MID_GREY_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )
    add_text_box(
        slide, data["date"],
        Inches(5.9), Inches(3.7), Inches(7), Inches(0.4),
        size=12, color=MID_GREY_PPT, align=PP_ALIGN.LEFT, font="Calibri"
    )

    # Decorative circle (light accent, bottom-right of right panel)
    circle = slide.shapes.add_shape(9, Inches(8.5), Inches(4.5), Inches(3.5), Inches(3.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PptRGB(0xE8, 0xCC, 0xD8)  # light burgundy tint
    circle.line.fill.background()
    return slide


def build_bullets_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")

    y = Inches(1.5)
    for label, body in data["bullets"]:
        # Label pill
        add_rect(slide, Inches(0.4), y, Inches(2.8), Inches(0.38), fill_color=ORANGE_PPT)
        add_text_box(slide, label, Inches(0.45), y + Inches(0.02), Inches(2.7), Inches(0.34),
                     bold=True, size=13, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
        # Body
        add_text_box(slide, body, Inches(0.4), y + Inches(0.44), SLIDE_W - Inches(0.9), Inches(0.7),
                     size=14, color=DARK_PPT, align=PP_ALIGN.LEFT, font="Calibri")
        y += Inches(1.35)

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_grid_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, color=MID_GREY_PPT, font="Calibri")

    items = data["items"]
    col_w = (SLIDE_W - Inches(1)) / len(items)
    for i, (emoji, title, desc) in enumerate(items):
        x = Inches(0.5) + i * col_w
        # Card
        add_rect(slide, x, Inches(1.85), col_w - Inches(0.15), Inches(4.5),
                 fill_color=LIGHT_BG_PPT)
        # Top accent
        add_rect(slide, x, Inches(1.85), col_w - Inches(0.15), Inches(0.12), fill_color=BURGUNDY_PPT)
        add_text_box(slide, emoji, x + Inches(0.1), Inches(2.1), col_w - Inches(0.25), Inches(0.7),
                     size=30, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, title, x + Inches(0.1), Inches(2.9), col_w - Inches(0.25), Inches(0.6),
                     bold=True, size=13, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, desc, x + Inches(0.1), Inches(3.6), col_w - Inches(0.25), Inches(1.2),
                     size=11, color=MID_GREY_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_flow_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, color=MID_GREY_PPT, font="Calibri")

    steps = data["steps"]
    step_w = (SLIDE_W - Inches(1)) / len(steps)
    arrow_w = Inches(0.4)
    box_w = step_w - arrow_w

    for i, (emoji, label, desc) in enumerate(steps):
        x = Inches(0.5) + i * step_w
        # Step box
        add_rect(slide, x, Inches(2.0), box_w, Inches(4.0), fill_color=LIGHT_BG_PPT)
        add_rect(slide, x, Inches(2.0), box_w, Inches(0.15), fill_color=ORANGE_PPT)
        add_text_box(slide, emoji, x, Inches(2.3), box_w, Inches(0.8),
                     size=32, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, label, x, Inches(3.2), box_w, Inches(0.5),
                     bold=True, size=14, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, desc, x + Inches(0.1), Inches(3.8), box_w - Inches(0.2), Inches(1.8),
                     size=11, color=DARK_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        # Arrow between boxes (not after last)
        if i < len(steps) - 1:
            ax = x + box_w + Inches(0.05)
            add_text_box(slide, "→", ax, Inches(3.0), arrow_w, Inches(0.6),
                         bold=True, size=20, color=ORANGE_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_four_cards_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, color=MID_GREY_PPT, font="Calibri")

    cards = data["cards"]
    card_w = (SLIDE_W - Inches(1)) / len(cards)
    for i, (emoji, title, desc) in enumerate(cards):
        x = Inches(0.5) + i * card_w
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.15), Inches(4.5), fill_color=LIGHT_BG_PPT)
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.15), Inches(0.15), fill_color=BURGUNDY_PPT)
        add_text_box(slide, emoji, x, Inches(2.3), card_w - Inches(0.15), Inches(0.8),
                     size=34, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, title, x, Inches(3.2), card_w - Inches(0.15), Inches(0.5),
                     bold=True, size=15, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, desc, x + Inches(0.1), Inches(3.85), card_w - Inches(0.35), Inches(2.0),
                     size=12, color=DARK_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_detail_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")

    add_text_box(slide, data["lead"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=14, italic=True, color=ORANGE_PPT, font="Calibri")

    y = Inches(1.8)
    for label, body in data["points"]:
        # Label
        add_rect(slide, Inches(0.4), y, Inches(3.8), Inches(0.36), fill_color=BURGUNDY_PPT)
        add_text_box(slide, label, Inches(0.5), y + Inches(0.02), Inches(3.6), Inches(0.32),
                     bold=True, size=12, color=WHITE_PPT, font="Calibri")
        # Body
        add_text_box(slide, body, Inches(4.4), y, SLIDE_W - Inches(5.0), Inches(0.55),
                     size=12, color=DARK_PPT, font="Calibri")
        y += Inches(0.75)

    if data.get("callout"):
        add_rect(slide, Inches(0.4), y + Inches(0.1), SLIDE_W - Inches(0.8), Inches(0.7), fill_color=ORANGE_PPT)
        add_text_box(slide, "💡  " + data["callout"],
                     Inches(0.6), y + Inches(0.15), SLIDE_W - Inches(1.2), Inches(0.6),
                     bold=True, size=12, color=WHITE_PPT, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_three_steps_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, italic=True, color=ORANGE_PPT, font="Calibri")

    steps = data["steps"]
    col_w = (SLIDE_W - Inches(1)) / len(steps)
    for i, (header, body) in enumerate(steps):
        x = Inches(0.5) + i * col_w
        # Column card
        add_rect(slide, x, Inches(1.9), col_w - Inches(0.2), Inches(5.0), fill_color=LIGHT_BG_PPT)
        # Top header bar
        add_rect(slide, x, Inches(1.9), col_w - Inches(0.2), Inches(0.9),
                 fill_color=BURGUNDY_PPT)
        add_text_box(slide, header, x + Inches(0.1), Inches(1.95), col_w - Inches(0.4), Inches(0.8),
                     bold=True, size=14, color=WHITE_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        # Body text
        add_text_box(slide, body, x + Inches(0.15), Inches(2.9), col_w - Inches(0.4), Inches(3.5),
                     size=13, color=DARK_PPT, align=PP_ALIGN.LEFT, font="Calibri")
        # Arrow connector (not after last)
        if i < len(steps) - 1:
            ax = x + col_w - Inches(0.25)
            add_text_box(slide, "→", ax, Inches(3.8), Inches(0.5), Inches(0.6),
                         bold=True, size=22, color=ORANGE_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_numbered_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, italic=True, color=ORANGE_PPT, font="Calibri")

    steps = data["steps"]
    col1_steps = steps[:3]
    col2_steps = steps[3:]

    for col_idx, col_steps in enumerate([col1_steps, col2_steps]):
        x = Inches(0.5) + col_idx * Inches(6.3)
        y = Inches(1.9)
        for j, step in enumerate(col_steps):
            num = col_idx * 3 + j + 1
            # Number circle
            add_rect(slide, x, y, Inches(0.5), Inches(0.5), fill_color=BURGUNDY_PPT)
            add_text_box(slide, str(num), x, y, Inches(0.5), Inches(0.5),
                         bold=True, size=14, color=WHITE_PPT, align=PP_ALIGN.CENTER, font="Calibri")
            # Step text
            add_text_box(slide, step, x + Inches(0.6), y, Inches(5.5), Inches(0.6),
                         size=13, color=DARK_PPT, font="Calibri")
            y += Inches(0.85)

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_reminders_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=13, italic=True, color=ORANGE_PPT, font="Calibri")

    reminders = data["reminders"]
    card_w = (SLIDE_W - Inches(1)) / len(reminders)
    for i, (emoji, title, trigger, recipient) in enumerate(reminders):
        x = Inches(0.5) + i * card_w
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.15), Inches(4.6), fill_color=LIGHT_BG_PPT)
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.15), Inches(0.15), fill_color=ORANGE_PPT)
        add_text_box(slide, emoji, x, Inches(2.25), card_w - Inches(0.15), Inches(0.7),
                     size=30, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, title, x + Inches(0.1), Inches(3.0), card_w - Inches(0.3), Inches(0.5),
                     bold=True, size=13, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, f"When:\n{trigger}", x + Inches(0.1), Inches(3.6), card_w - Inches(0.3), Inches(1.2),
                     size=11, color=DARK_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, f"To: {recipient}", x + Inches(0.1), Inches(4.9), card_w - Inches(0.3), Inches(0.5),
                     bold=True, size=11, color=ORANGE_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_getting_started_slide(prs, data):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_color=BURGUNDY_PPT)
    add_text_box(slide, data["title"], Inches(0.4), Inches(0.2), SLIDE_W - Inches(0.8), Inches(0.8),
                 bold=True, size=28, color=WHITE_PPT, align=PP_ALIGN.LEFT, font="Calibri")
    add_text_box(slide, data["intro"], Inches(0.4), Inches(1.25), SLIDE_W - Inches(0.8), Inches(0.4),
                 size=14, italic=True, color=ORANGE_PPT, font="Calibri")

    actions = data["actions"]
    card_w = (SLIDE_W - Inches(1)) / len(actions)
    for i, (emoji, title, desc) in enumerate(actions):
        x = Inches(0.5) + i * card_w
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.2), Inches(3.8), fill_color=LIGHT_BG_PPT)
        add_rect(slide, x, Inches(2.0), card_w - Inches(0.2), Inches(0.15), fill_color=BURGUNDY_PPT)
        add_text_box(slide, emoji, x, Inches(2.3), card_w - Inches(0.2), Inches(0.8),
                     size=36, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, title, x + Inches(0.1), Inches(3.2), card_w - Inches(0.4), Inches(0.5),
                     bold=True, size=15, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
        add_text_box(slide, desc, x + Inches(0.1), Inches(3.8), card_w - Inches(0.4), Inches(1.5),
                     size=12, color=DARK_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_text_box(slide, f"💬  {data['note']}",
                 Inches(0.5), Inches(6.0), SLIDE_W - Inches(1), Inches(0.6),
                 size=13, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, bold=True, font="Calibri")

    add_footer(slide, "AktivAsia Portal  |  User Guide & SOP  |  May 2026")
    return slide


def build_contact_slide(prs, data):
    slide = blank_slide(prs)
    # Full burgundy background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=BURGUNDY_PPT)

    add_text_box(slide, data["title"],
                 Inches(0.4), Inches(0.8), SLIDE_W - Inches(0.8), Inches(0.9),
                 bold=True, size=32, color=WHITE_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    # White card in center
    add_rect(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(3.5), fill_color=WHITE_PPT)
    add_rect(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(0.15), fill_color=ORANGE_PPT)

    add_text_box(slide, data["name"],
                 Inches(2.7), Inches(2.3), Inches(7.9), Inches(0.7),
                 bold=True, size=22, color=BURGUNDY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
    add_text_box(slide, data["role"],
                 Inches(2.7), Inches(3.05), Inches(7.9), Inches(0.5),
                 size=14, color=MID_GREY_PPT, align=PP_ALIGN.CENTER, font="Calibri")
    add_text_box(slide, f"📧  {data['email']}",
                 Inches(2.7), Inches(3.65), Inches(7.9), Inches(0.5),
                 size=14, color=ORANGE_PPT, align=PP_ALIGN.CENTER, font="Calibri")

    add_text_box(slide, data["tagline"],
                 Inches(0.4), Inches(6.0), SLIDE_W - Inches(0.8), Inches(0.6),
                 size=16, color=WHITE_PPT, align=PP_ALIGN.CENTER, italic=True, font="Calibri")


def add_screenshot_pptx(slide, img_path: str, left, top, width, height):
    """Insert a screenshot image onto a PPT slide."""
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(img_path, left, top, width, height)


SLIDE_BUILDERS = {
    "cover":         build_cover,
    "bullets":       build_bullets_slide,
    "grid":          build_grid_slide,
    "flow":          build_flow_slide,
    "four_cards":    build_four_cards_slide,
    "detail":        build_detail_slide,
    "three_steps":   build_three_steps_slide,
    "numbered":      build_numbered_slide,
    "reminders":     build_reminders_slide,
    "getting_started": build_getting_started_slide,
    "contact":       build_contact_slide,
}


def build_pptx():
    prs = new_prs()
    for slide_data in SLIDES:
        builder = SLIDE_BUILDERS.get(slide_data["type"])
        if builder:
            slide = builder(prs, slide_data)
            # If this slide has a screenshot and builder returned the slide object,
            # add a framed screenshot thumbnail in the bottom-right corner
            img_name = slide_data.get("screenshot")
            if img_name and slide is not None:
                img_path = screen(img_name)
                if img_path:
                    # White frame
                    add_rect(slide,
                             SLIDE_W - Inches(5.1), SLIDE_H - Inches(3.0),
                             Inches(4.8), Inches(2.7),
                             fill_color=WHITE_PPT)
                    # Screenshot inside frame
                    add_screenshot_pptx(slide, img_path,
                                        SLIDE_W - Inches(5.0), SLIDE_H - Inches(2.9),
                                        Inches(4.6), Inches(2.6))
                    # Label
                    add_text_box(slide, "📸  Portal screenshot",
                                 SLIDE_W - Inches(5.1), SLIDE_H - Inches(0.72),
                                 Inches(4.8), Inches(0.25),
                                 size=8, color=MID_GREY_PPT,
                                 align=PP_ALIGN.CENTER, font="Calibri")
        else:
            print(f"⚠️  Unknown slide type: {slide_data['type']}")

    out = OUT_DIR / "AktivAsia_Portal_Presentation.pptx"
    prs.save(str(out))
    print(f"✅  Saved: {out}")


# ─── ENTRY POINT ──────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import sys
    if sys.stdout.encoding and sys.stdout.encoding.lower() != "utf-8":
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    print("Building AktivAsia Portal documentation...")
    build_docx()
    build_pptx()
    print("\nDone! Both files are in the docs/ folder.")
